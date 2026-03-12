from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(__file__).with_name("native_elements_demo.pptx")


def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.5), Inches(0.6))
    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(26)
    paragraph.font.bold = True


def build_demo_deck():
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Native Elements Demo")

    subtitle = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(4.8), Inches(0.8))
    p = subtitle.text_frame.paragraphs[0]
    p.text = "This slide contains editable text, vector shapes, a live table, and a live chart."
    p.font.name = "Aptos"
    p.font.size = Pt(18)

    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(1.8),
        Inches(3.4),
        Inches(0.8),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(14, 116, 144)
    banner.line.color.rgb = RGBColor(7, 54, 66)
    banner.text_frame.text = "Status: Native PowerPoint Shape"
    banner.text_frame.paragraphs[0].font.size = Pt(16)
    banner.text_frame.paragraphs[0].font.bold = True
    banner.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    callout = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CLOUD_CALLOUT,
        Inches(4.3),
        Inches(1.65),
        Inches(2.1),
        Inches(1.2),
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(255, 243, 205)
    callout.line.color.rgb = RGBColor(184, 134, 11)
    callout.text_frame.text = "Should stay editable"
    callout.text_frame.paragraphs[0].font.size = Pt(14)

    table_shape = slide.shapes.add_table(3, 3, Inches(0.6), Inches(3.0), Inches(4.2), Inches(1.7))
    table = table_shape.table
    headers = ["Owner", "Workstream", "Progress"]
    rows = [
        ("Design", "Layout parsing", "80%"),
        ("Data", "Chart rebuild", "55%"),
    ]

    for col, text in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(217, 237, 247)

    for row_index, row_values in enumerate(rows, start=1):
        for col_index, text in enumerate(row_values):
            table.cell(row_index, col_index).text = text

    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("Detected", (4, 7, 9, 11))
    chart_data.add_series("Editable", (2, 5, 7, 9))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(5.1),
        Inches(3.0),
        Inches(4.2),
        Inches(2.8),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Rebuild Progress"

    notes = slide.shapes.add_textbox(Inches(0.7), Inches(5.1), Inches(4.0), Inches(0.8))
    notes_tf = notes.text_frame
    notes_tf.text = "Expected extractor report:"
    notes_tf.paragraphs[0].font.bold = True
    p2 = notes_tf.add_paragraph()
    p2.text = "1 text box, 2 shapes, 1 table, 1 chart"
    p2.font.size = Pt(14)

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    output = build_demo_deck()
    print(output)
