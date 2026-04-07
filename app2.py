import io
import os
from pathlib import Path

import cv2
import easyocr
import numpy as np
import streamlit as st
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

try:
    import fitz
except ImportError:  # pragma: no cover - handled in the UI/runtime path
    fitz = None


OCR_LANGUAGES = ["en"]
MIN_OCR_CONFIDENCE = 0.2
DEFAULT_FONT_SIZE_PT = 14
OCR_RENDER_SCALE = 2.0
DEFAULT_IMAGE_WIDTH_PTS = 720
MAX_SLIDE_DIMENSION_PTS = 4032
MAX_OCR_IMAGE_SIDE_PX = 2200
MIN_FONT_SIZE_PT = 10
MAX_FONT_SIZE_PT = 30


def _easyocr_cache_dir() -> str:
    # Streamlit deployments can fail if EasyOCR tries to write under HOME.
    cache_dir = Path(os.getenv("EASYOCR_MODULE_STORAGE", "/tmp/easyocr"))
    cache_dir.mkdir(parents=True, exist_ok=True)
    return str(cache_dir)


@st.cache_resource(show_spinner=False)
def load_ocr():
    return easyocr.Reader(
        OCR_LANGUAGES,
        model_storage_directory=_easyocr_cache_dir(),
    )


def _shape_bounds_to_points(shape):
    return (
        shape.left.pt,
        shape.top.pt,
        shape.width.pt,
        shape.height.pt,
    )


def _rect_from_xywh(left, top, width, height):
    return (left, top, left + width, top + height)


def _rect_intersection(rect_a, rect_b):
    left = max(rect_a[0], rect_b[0])
    top = max(rect_a[1], rect_b[1])
    right = min(rect_a[2], rect_b[2])
    bottom = min(rect_a[3], rect_b[3])
    if right <= left or bottom <= top:
        return None
    return (left, top, right, bottom)


def _rect_contains(outer_rect, inner_rect):
    return (
        inner_rect[0] >= outer_rect[0]
        and inner_rect[1] >= outer_rect[1]
        and inner_rect[2] <= outer_rect[2]
        and inner_rect[3] <= outer_rect[3]
    )


def _sanitize_crop_rect(rect, max_width, max_height):
    left, top, right, bottom = rect
    left = int(max(0, min(round(left), max_width - 1)))
    top = int(max(0, min(round(top), max_height - 1)))
    right = int(max(0, min(round(right), max_width)))
    bottom = int(max(0, min(round(bottom), max_height)))
    if right <= left or bottom <= top:
        return None
    return left, top, right, bottom


def _selection_bounds(width, height, selection):
    if not selection or not selection.get("enabled"):
        return 0, 0, width, height

    x = (selection["x_pct"] / 100.0) * width
    y = (selection["y_pct"] / 100.0) * height
    selected_width = (selection["width_pct"] / 100.0) * width
    selected_height = (selection["height_pct"] / 100.0) * height

    selected_width = min(max(selected_width, 1), max(width - x, 1))
    selected_height = min(max(selected_height, 1), max(height - y, 1))
    return x, y, selected_width, selected_height


def _shape_bounds_object(left_pt, top_pt, width_pt, height_pt):
    return type(
        "ShapeBounds",
        (),
        {
            "left": Pt(left_pt),
            "top": Pt(top_pt),
            "width": Pt(width_pt),
            "height": Pt(height_pt),
        },
    )()


def _add_picture_from_blob(new_slide, image_blob, left, top, width, height):
    new_slide.shapes.add_picture(
        io.BytesIO(image_blob),
        Pt(left) if isinstance(left, (int, float)) else left,
        Pt(top) if isinstance(top, (int, float)) else top,
        width=Pt(width) if isinstance(width, (int, float)) else width,
        height=Pt(height) if isinstance(height, (int, float)) else height,
    )


def _copy_color(source_color, target_color):
    try:
        if source_color.rgb is not None:
            target_color.rgb = RGBColor(*source_color.rgb)
    except Exception:
        return


def _copy_fill(source_fill, target_fill):
    try:
        if source_fill.type is None:
            return
        if source_fill.fore_color.rgb is not None:
            target_fill.solid()
            _copy_color(source_fill.fore_color, target_fill.fore_color)
    except Exception:
        return


def _copy_line(source_line, target_line):
    try:
        if source_line.color.rgb is not None:
            _copy_color(source_line.color, target_line.color)
        if source_line.width is not None:
            target_line.width = source_line.width
    except Exception:
        return


def _copy_font(source_font, target_font):
    for attr in ("bold", "italic", "underline", "name", "size"):
        try:
            value = getattr(source_font, attr)
        except Exception:
            continue
        if value is not None:
            setattr(target_font, attr, value)

    try:
        if source_font.color.rgb is not None:
            _copy_color(source_font.color, target_font.color)
    except Exception:
        return


def _copy_text_frame(source_text_frame, target_text_frame):
    target_text_frame.clear()
    target_text_frame.word_wrap = source_text_frame.word_wrap

    for index, source_paragraph in enumerate(source_text_frame.paragraphs):
        paragraph = (
            target_text_frame.paragraphs[0]
            if index == 0
            else target_text_frame.add_paragraph()
        )
        paragraph.alignment = source_paragraph.alignment
        paragraph.level = source_paragraph.level

        if source_paragraph.runs:
            for source_run in source_paragraph.runs:
                run = paragraph.add_run()
                run.text = source_run.text
                _copy_font(source_run.font, run.font)
        else:
            paragraph.text = source_paragraph.text
            _copy_font(source_paragraph.font, paragraph.font)


def _clone_text_like_shape(shape, new_slide, offset_left_pt=0, offset_top_pt=0):
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        cloned_shape = new_slide.shapes.add_shape(
            shape.auto_shape_type,
            Pt(shape.left.pt - offset_left_pt),
            Pt(shape.top.pt - offset_top_pt),
            shape.width,
            shape.height,
        )
        _copy_fill(shape.fill, cloned_shape.fill)
        _copy_line(shape.line, cloned_shape.line)
    else:
        cloned_shape = new_slide.shapes.add_textbox(
            Pt(shape.left.pt - offset_left_pt),
            Pt(shape.top.pt - offset_top_pt),
            shape.width,
            shape.height,
        )

    if shape.has_text_frame:
        _copy_text_frame(shape.text_frame, cloned_shape.text_frame)
    return True


def _clone_table_shape(shape, new_slide, offset_left_pt=0, offset_top_pt=0):
    source_table = shape.table
    rows = len(source_table.rows)
    cols = len(source_table.columns)
    table_shape = new_slide.shapes.add_table(
        rows,
        cols,
        Pt(shape.left.pt - offset_left_pt),
        Pt(shape.top.pt - offset_top_pt),
        shape.width,
        shape.height,
    )
    target_table = table_shape.table

    for col_index, column in enumerate(source_table.columns):
        target_table.columns[col_index].width = column.width

    for row_index, row in enumerate(source_table.rows):
        target_table.rows[row_index].height = row.height

    for row_index in range(rows):
        for col_index in range(cols):
            source_cell = source_table.cell(row_index, col_index)
            target_cell = target_table.cell(row_index, col_index)
            target_cell.text = source_cell.text
            _copy_fill(source_cell.fill, target_cell.fill)
            _copy_text_frame(source_cell.text_frame, target_cell.text_frame)

    return True


def _clone_category_chart(shape, new_slide, offset_left_pt=0, offset_top_pt=0):
    chart = shape.chart
    plot = chart.plots[0]
    categories = [category.label for category in plot.categories]
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for series in chart.series:
        chart_data.add_series(series.name, tuple(series.values))

    chart_shape = new_slide.shapes.add_chart(
        chart.chart_type,
        Pt(shape.left.pt - offset_left_pt),
        Pt(shape.top.pt - offset_top_pt),
        shape.width,
        shape.height,
        chart_data,
    )
    target_chart = chart_shape.chart

    try:
        if chart.has_title:
            target_chart.has_title = True
            target_chart.chart_title.text_frame.text = chart.chart_title.text_frame.text
    except Exception:
        pass

    return True


def _clone_chart_shape(shape, new_slide, offset_left_pt=0, offset_top_pt=0):
    try:
        return _clone_category_chart(shape, new_slide, offset_left_pt, offset_top_pt)
    except Exception:
        placeholder = new_slide.shapes.add_textbox(
            Pt(shape.left.pt - offset_left_pt),
            Pt(shape.top.pt - offset_top_pt),
            shape.width,
            shape.height,
        )
        placeholder.text_frame.text = (
            "Chart detected. Auto-rebuild is not yet supported for this chart type."
        )
        return False


def _clone_native_shape(shape, new_slide, offset_left_pt=0, offset_top_pt=0):
    if getattr(shape, "has_table", False):
        return _clone_table_shape(shape, new_slide, offset_left_pt, offset_top_pt), "tables"

    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        cloned = _clone_chart_shape(shape, new_slide, offset_left_pt, offset_top_pt)
        return cloned, "charts" if cloned else "unsupported"

    if shape.has_text_frame or shape.shape_type in {
        MSO_SHAPE_TYPE.AUTO_SHAPE,
        MSO_SHAPE_TYPE.CALLOUT,
        MSO_SHAPE_TYPE.FREEFORM,
        MSO_SHAPE_TYPE.PLACEHOLDER,
        MSO_SHAPE_TYPE.TEXT_BOX,
    }:
        return _clone_text_like_shape(shape, new_slide, offset_left_pt, offset_top_pt), "text_shapes"

    return False, "unsupported"


def _clone_shape_as_is(shape, new_slide):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        _add_picture_from_blob(
            new_slide,
            shape.image.blob,
            shape.left,
            shape.top,
            shape.width,
            shape.height,
        )
        return True, "pictures"

    return _clone_native_shape(shape, new_slide)


def _extract_text_blocks(reader, image_np: np.ndarray):
    results = reader.readtext(image_np)
    img_h, img_w = image_np.shape[:2]
    mask = np.zeros((img_h, img_w), dtype=np.uint8)
    blocks = []

    for bbox, text, prob in results:
        if prob < MIN_OCR_CONFIDENCE:
            continue

        points = np.array(bbox, dtype=np.int32)
        if points.shape != (4, 2):
            continue

        cv2.fillPoly(mask, [points], 255)
        x_coords = points[:, 0]
        y_coords = points[:, 1]
        left = int(x_coords.min())
        top = int(y_coords.min())
        right = int(x_coords.max())
        bottom = int(y_coords.max())

        if right <= left or bottom <= top:
            continue

        blocks.append(
            {
                "text": text,
                "left_px": left,
                "top_px": top,
                "width_px": right - left,
                "height_px": bottom - top,
                "right_px": right,
                "bottom_px": bottom,
                "confidence": float(prob),
            }
        )

    return blocks, mask


def _looks_like_bullet(text):
    stripped = text.lstrip()
    return stripped.startswith(("-", "*", "\u2022", "\u25e6", "\u25aa"))


def _estimate_font_size_from_line_height(line_height_px, img_h, shape_height_pts):
    line_height_pts = (line_height_px / max(img_h, 1)) * shape_height_pts
    estimated = line_height_pts * 0.8
    return max(MIN_FONT_SIZE_PT, min(estimated, MAX_FONT_SIZE_PT))


def _group_ocr_blocks(blocks):
    if not blocks:
        return []

    ordered_blocks = sorted(blocks, key=lambda block: (block["top_px"], block["left_px"]))
    lines = []

    for block in ordered_blocks:
        block_center = block["top_px"] + (block["height_px"] / 2)
        matched_line = None

        for line in lines:
            line_center = line["top_px"] + (line["height_px"] / 2)
            vertical_threshold = max(line["height_px"], block["height_px"]) * 0.65
            if abs(block_center - line_center) <= vertical_threshold:
                matched_line = line
                break

        if matched_line is None:
            lines.append(
                {
                    "items": [block],
                    "left_px": block["left_px"],
                    "top_px": block["top_px"],
                    "right_px": block["left_px"] + block["width_px"],
                    "bottom_px": block["top_px"] + block["height_px"],
                    "height_px": block["height_px"],
                }
            )
            continue

        matched_line["items"].append(block)
        matched_line["left_px"] = min(matched_line["left_px"], block["left_px"])
        matched_line["top_px"] = min(matched_line["top_px"], block["top_px"])
        matched_line["right_px"] = max(
            matched_line["right_px"], block["left_px"] + block["width_px"]
        )
        matched_line["bottom_px"] = max(
            matched_line["bottom_px"], block["top_px"] + block["height_px"]
        )
        matched_line["height_px"] = matched_line["bottom_px"] - matched_line["top_px"]

    normalized_lines = []
    for line in lines:
        items = sorted(line["items"], key=lambda item: item["left_px"])
        text = " ".join(item["text"] for item in items).strip()
        left_px = line["left_px"]
        top_px = line["top_px"]
        right_px = line["right_px"]
        bottom_px = line["bottom_px"]
        normalized_lines.append(
            {
                "text": text,
                "left_px": left_px,
                "top_px": top_px,
                "right_px": right_px,
                "bottom_px": bottom_px,
                "height_px": line["height_px"],
                "bullet_like": _looks_like_bullet(text),
            }
        )

    normalized_lines.sort(key=lambda line: (line["top_px"], line["left_px"]))
    groups = []

    for line in normalized_lines:
        if not line["text"]:
            continue

        if not groups:
            groups.append(
                {
                    "text": line["text"],
                    "left_px": line["left_px"],
                    "top_px": line["top_px"],
                    "right_px": line["right_px"],
                    "bottom_px": line["bottom_px"],
                    "line_height_px": line["height_px"],
                    "lines": [line],
                    "bullet_like": line["bullet_like"],
                }
            )
            continue

        current = groups[-1]
        vertical_gap = line["top_px"] - current["bottom_px"]
        line_height_threshold = max(current["line_height_px"], line["height_px"]) * 2.2
        left_delta = abs(line["left_px"] - current["left_px"])
        max_left_delta = max(current["right_px"] - current["left_px"], 1) * 0.4
        indent_shift = line["left_px"] - current["left_px"]
        current_font_height = current["line_height_px"]
        line_height_ratio = max(current_font_height, line["height_px"]) / max(
            min(current_font_height, line["height_px"]),
            1,
        )
        current_is_title_like = current_font_height >= line["height_px"] * 1.35 and len(
            current["lines"]
        ) == 1
        different_structure = current["bullet_like"] != line["bullet_like"]
        should_start_new_group = (
            vertical_gap > line_height_threshold
            or left_delta > max_left_delta
            or indent_shift > max(current_font_height, line["height_px"]) * 0.75
            or line_height_ratio > 1.45
            or (current_is_title_like and not current["bullet_like"] and not line["bullet_like"])
            or different_structure
        )

        if not should_start_new_group:
            current["text"] = f"{current['text']}\n{line['text']}"
            current["left_px"] = min(current["left_px"], line["left_px"])
            current["top_px"] = min(current["top_px"], line["top_px"])
            current["right_px"] = max(current["right_px"], line["right_px"])
            current["bottom_px"] = max(current["bottom_px"], line["bottom_px"])
            current["line_height_px"] = max(current["line_height_px"], line["height_px"])
            current["lines"].append(line)
        else:
            groups.append(
                {
                    "text": line["text"],
                    "left_px": line["left_px"],
                    "top_px": line["top_px"],
                    "right_px": line["right_px"],
                    "bottom_px": line["bottom_px"],
                    "line_height_px": line["height_px"],
                    "lines": [line],
                    "bullet_like": line["bullet_like"],
                }
            )

    return [
        {
            "text": group["text"],
            "left_px": group["left_px"],
            "top_px": group["top_px"],
            "width_px": group["right_px"] - group["left_px"],
            "height_px": group["bottom_px"] - group["top_px"],
            "lines": group["lines"],
        }
        for group in groups
    ]


def _clean_image(image_np: np.ndarray, mask: np.ndarray) -> bytes:
    if np.count_nonzero(mask) == 0:
        cleaned_img_np = image_np
    else:
        cleaned_img_np = cv2.inpaint(image_np, mask, 3, cv2.INPAINT_TELEA)

    cleaned_pil = Image.fromarray(cleaned_img_np)
    img_bytes = io.BytesIO()
    cleaned_pil.save(img_bytes, format="PNG")
    img_bytes.seek(0)
    return img_bytes.getvalue()


def _resize_image_for_ocr(image_np: np.ndarray, max_side_px=MAX_OCR_IMAGE_SIDE_PX):
    img_h, img_w = image_np.shape[:2]
    longest_side = max(img_h, img_w)
    if longest_side <= max_side_px:
        return image_np

    scale = max_side_px / longest_side
    resized = cv2.resize(
        image_np,
        (max(int(img_w * scale), 1), max(int(img_h * scale), 1)),
        interpolation=cv2.INTER_AREA,
    )
    return resized


def _add_text_boxes(new_slide, blocks, img_w, img_h, shape):
    shape_left, shape_top, shape_width, shape_height = _shape_bounds_to_points(shape)

    for block in blocks:
        left = shape_left + (block["left_px"] / img_w) * shape_width
        top = shape_top + (block["top_px"] / img_h) * shape_height
        width = (block["width_px"] / img_w) * shape_width
        height = (block["height_px"] / img_h) * shape_height

        textbox = new_slide.shapes.add_textbox(
            Pt(left),
            Pt(top),
            Pt(max(width, 1)),
            Pt(max(height, 1)),
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        lines = block.get("lines") or [{"text": block["text"], "height_px": block["height_px"]}]

        for index, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = line["text"]
            paragraph.font.size = Pt(
                _estimate_font_size_from_line_height(line["height_px"], img_h, shape_height)
            )


def _build_output_presentation(width_pts=720, height_pts=540):
    prs = Presentation()
    prs.slide_width = Pt(width_pts)
    prs.slide_height = Pt(height_pts)
    return prs


def _save_presentation(output_prs):
    output_stream = io.BytesIO()
    output_prs.save(output_stream)
    output_stream.seek(0)
    return output_stream.getvalue()


def _full_slide_shape(output_prs):
    return _shape_bounds_object(0, 0, output_prs.slide_width.pt, output_prs.slide_height.pt)


def _fit_image_to_slide(image_width_px, image_height_px, target_width_pts=DEFAULT_IMAGE_WIDTH_PTS):
    height_pts = target_width_pts * (image_height_px / max(image_width_px, 1))
    width_pts = target_width_pts

    if height_pts > MAX_SLIDE_DIMENSION_PTS:
        scale = MAX_SLIDE_DIMENSION_PTS / height_pts
        width_pts *= scale
        height_pts = MAX_SLIDE_DIMENSION_PTS

    if width_pts > MAX_SLIDE_DIMENSION_PTS:
        scale = MAX_SLIDE_DIMENSION_PTS / width_pts
        width_pts = MAX_SLIDE_DIMENSION_PTS
        height_pts *= scale

    return max(width_pts, 1), max(height_pts, 1)


def _crop_array_to_rect(image_np, rect):
    safe_rect = _sanitize_crop_rect(rect, image_np.shape[1], image_np.shape[0])
    if safe_rect is None:
        return None
    left, top, right, bottom = safe_rect
    cropped = image_np[top:bottom, left:right]
    return cropped if cropped.size else None


def _crop_pil_bytes_to_rect(image_bytes, rect_px):
    image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    safe_rect = _sanitize_crop_rect(rect_px, image.size[0], image.size[1])
    if safe_rect is None:
        return None
    cropped = image.crop(safe_rect)
    buffer = io.BytesIO()
    cropped.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer.getvalue()


def _overlay_editable_region(new_slide, image_np, target_shape, report, report_key_prefix):
    if image_np is None or image_np.size == 0:
        report["unsupported"] += 1
        return

    reader = load_ocr()
    image_np = _resize_image_for_ocr(image_np)
    img_h, img_w = image_np.shape[:2]
    blocks, mask = _extract_text_blocks(reader, image_np)
    grouped_blocks = _group_ocr_blocks(blocks)
    cleaned_image_bytes = _clean_image(image_np, mask)

    _add_picture_from_blob(
        new_slide,
        cleaned_image_bytes,
        target_shape.left,
        target_shape.top,
        target_shape.width,
        target_shape.height,
    )
    _add_text_boxes(
        new_slide,
        grouped_blocks,
        img_w,
        img_h,
        target_shape,
    )

    report[f"{report_key_prefix}_ocr_regions"] += len(blocks)
    report[f"{report_key_prefix}_text_groups"] += len(grouped_blocks)


def _process_flat_image(image_np, output_prs, report, report_key_prefix):
    new_slide = output_prs.slides.add_slide(output_prs.slide_layouts[6])
    _overlay_editable_region(new_slide, image_np, _full_slide_shape(output_prs), report, report_key_prefix)


def process_pptx_advanced(input_file, selection=None):
    prs = Presentation(input_file)
    if selection and selection.get("enabled"):
        slide = prs.slides[selection["page_number"] - 1]
        selection_x, selection_y, selection_width, selection_height = _selection_bounds(
            prs.slide_width.pt, prs.slide_height.pt, selection
        )
        output_prs = _build_output_presentation(prs.slide_width.pt, prs.slide_height.pt)
        slides_to_process = [
            (
                slide,
                selection_x,
                selection_y,
                selection_width,
                selection_height,
                True,
            )
        ]
    else:
        output_prs = _build_output_presentation(prs.slide_width.pt, prs.slide_height.pt)
        slides_to_process = [
            (slide, 0, 0, prs.slide_width.pt, prs.slide_height.pt, False) for slide in prs.slides
        ]

    reader = load_ocr()
    report = {
        "slides_processed": len(slides_to_process),
        "pictures_ocr_regions": 0,
        "pictures_text_groups": 0,
        "text_shapes": 0,
        "tables": 0,
        "charts": 0,
        "unsupported": 0,
    }

    for slide, selection_x, selection_y, selection_width, selection_height, selection_enabled in slides_to_process:
        new_slide = output_prs.slides.add_slide(output_prs.slide_layouts[6])
        selection_rect = _rect_from_xywh(selection_x, selection_y, selection_width, selection_height)

        if selection_enabled:
            for shape in slide.shapes:
                cloned, bucket = _clone_shape_as_is(shape, new_slide)
                if not cloned:
                    report["unsupported"] += 1
                elif bucket in report:
                    report[bucket] += 1

        for shape in slide.shapes:
            shape_left, shape_top, shape_width, shape_height = _shape_bounds_to_points(shape)
            shape_rect = _rect_from_xywh(shape_left, shape_top, shape_width, shape_height)
            intersection = _rect_intersection(shape_rect, selection_rect)
            if intersection is None:
                continue

            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                if selection_enabled:
                    continue
                if not _rect_contains(selection_rect, shape_rect):
                    report["unsupported"] += 1
                    continue

                _, bucket = _clone_native_shape(shape, new_slide, selection_x, selection_y)
                report[bucket] += 1
                continue

            image_stream = io.BytesIO(shape.image.blob)
            image = Image.open(image_stream).convert("RGB")
            image_np = np.array(image)
            crop_left = ((intersection[0] - shape_left) / max(shape_width, 1)) * image_np.shape[1]
            crop_top = ((intersection[1] - shape_top) / max(shape_height, 1)) * image_np.shape[0]
            crop_right = ((intersection[2] - shape_left) / max(shape_width, 1)) * image_np.shape[1]
            crop_bottom = ((intersection[3] - shape_top) / max(shape_height, 1)) * image_np.shape[0]
            image_np = _crop_array_to_rect(image_np, (crop_left, crop_top, crop_right, crop_bottom))
            if image_np is None:
                report["unsupported"] += 1
                continue
            target_shape = _shape_bounds_object(
                intersection[0] if selection_enabled else intersection[0] - selection_x,
                intersection[1] if selection_enabled else intersection[1] - selection_y,
                intersection[2] - intersection[0],
                intersection[3] - intersection[1],
            )
            _overlay_editable_region(
                new_slide,
                image_np,
                target_shape,
                report,
                "pictures",
            )

    return _save_presentation(output_prs), report


def _pdf_font_flags(span):
    font_name = span.get("font", "") or ""
    font_lower = font_name.lower()
    return {
        "bold": "bold" in font_lower,
        "italic": "italic" in font_lower or "oblique" in font_lower,
    }


def _pdf_text_blocks(page_dict):
    return [block for block in page_dict.get("blocks", []) if block.get("type") == 0]


def _pdf_image_blocks(page_dict):
    return [block for block in page_dict.get("blocks", []) if block.get("type") == 1]


def _add_pdf_text_block(new_slide, block, offset_x=0, offset_y=0, selection_rect=None):
    x0, y0, x1, y1 = block["bbox"]
    if selection_rect is not None:
        clipped = _rect_intersection((x0, y0, x1, y1), selection_rect)
        if clipped is None:
            return False
        x0, y0, x1, y1 = clipped
    width = max(x1 - x0, 1)
    height = max(y1 - y0, 1)
    textbox = new_slide.shapes.add_textbox(Pt(x0 - offset_x), Pt(y0 - offset_y), Pt(width), Pt(height))
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph_index = 0

    for line in block.get("lines", []):
        paragraph = (
            text_frame.paragraphs[0]
            if paragraph_index == 0
            else text_frame.add_paragraph()
        )
        paragraph_index += 1

        for span in line.get("spans", []):
            text = span.get("text", "")
            if not text:
                continue

            run = paragraph.add_run()
            run.text = text
            run.font.size = Pt(max(span.get("size", DEFAULT_FONT_SIZE_PT), 1))
            if span.get("font"):
                run.font.name = span["font"]

            flags = _pdf_font_flags(span)
            run.font.bold = flags["bold"]
            run.font.italic = flags["italic"]

    if not any(paragraph.text for paragraph in text_frame.paragraphs):
        new_slide.shapes._spTree.remove(textbox._element)
        return False

    return True


def _add_pdf_image_block(new_slide, block, offset_x=0, offset_y=0, selection_rect=None):
    image_bytes = block.get("image")
    if not image_bytes:
        return False

    x0, y0, x1, y1 = block["bbox"]
    if selection_rect is not None:
        clipped = _rect_intersection((x0, y0, x1, y1), selection_rect)
        if clipped is None:
            return False
        image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        crop_left = ((clipped[0] - x0) / max(x1 - x0, 1)) * image.size[0]
        crop_top = ((clipped[1] - y0) / max(y1 - y0, 1)) * image.size[1]
        crop_right = ((clipped[2] - x0) / max(x1 - x0, 1)) * image.size[0]
        crop_bottom = ((clipped[3] - y0) / max(y1 - y0, 1)) * image.size[1]
        image_bytes = _crop_pil_bytes_to_rect(image_bytes, (crop_left, crop_top, crop_right, crop_bottom))
        if image_bytes is None:
            return False
        x0, y0, x1, y1 = clipped
    width = max(x1 - x0, 1)
    height = max(y1 - y0, 1)
    new_slide.shapes.add_picture(
        io.BytesIO(image_bytes),
        Pt(x0 - offset_x),
        Pt(y0 - offset_y),
        width=Pt(width),
        height=Pt(height),
    )
    return True


def _render_pdf_page_for_ocr(page):
    matrix = fitz.Matrix(OCR_RENDER_SCALE, OCR_RENDER_SCALE)
    pix = page.get_pixmap(matrix=matrix, alpha=False)
    image = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
    return _resize_image_for_ocr(np.array(image))


def _render_pdf_page_image(page):
    pix = page.get_pixmap(alpha=False)
    return pix.tobytes("png")


def _ocr_pdf_page(page, new_slide, reader, selection_rect=None, offset_x=0, offset_y=0):
    page_image = _render_pdf_page_for_ocr(page)
    page_rect = _rect_from_xywh(0, 0, page.rect.width, page.rect.height)
    if selection_rect is not None:
        scale_x = page_image.shape[1] / max(page.rect.width, 1)
        scale_y = page_image.shape[0] / max(page.rect.height, 1)
        crop_rect = (
            selection_rect[0] * scale_x,
            selection_rect[1] * scale_y,
            selection_rect[2] * scale_x,
            selection_rect[3] * scale_y,
        )
        page_image = _crop_array_to_rect(page_image, crop_rect)
        page_rect = selection_rect
    img_h, img_w = page_image.shape[:2]
    blocks, mask = _extract_text_blocks(reader, page_image)
    grouped_blocks = _group_ocr_blocks(blocks)
    cleaned_image_bytes = _clean_image(page_image, mask)

    new_slide.shapes.add_picture(
        io.BytesIO(cleaned_image_bytes),
        Pt(0),
        Pt(0),
        width=Pt(page_rect[2] - page_rect[0]),
        height=Pt(page_rect[3] - page_rect[1]),
    )

    for block in grouped_blocks:
        left = (block["left_px"] / img_w) * (page_rect[2] - page_rect[0])
        top = (block["top_px"] / img_h) * (page_rect[3] - page_rect[1])
        width = (block["width_px"] / img_w) * (page_rect[2] - page_rect[0])
        height = (block["height_px"] / img_h) * (page_rect[3] - page_rect[1])

        textbox = new_slide.shapes.add_textbox(
            Pt(left),
            Pt(top),
            Pt(max(width, 1)),
            Pt(max(height, 1)),
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        for index, line in enumerate(block.get("lines", [])):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = line["text"]
            paragraph.font.size = Pt(
                _estimate_font_size_from_line_height(
                    line["height_px"], img_h, page_rect[3] - page_rect[1]
                )
            )

    return len(blocks), len(grouped_blocks)


def process_pdf_advanced(input_file, selection=None):
    if fitz is None:
        raise ImportError("PyMuPDF is required for PDF support. Add pymupdf to the app dependencies.")

    pdf_bytes = input_file.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    if doc.page_count == 0:
        raise ValueError("The uploaded PDF has no pages.")

    if selection and selection.get("enabled"):
        pages_to_process = [doc[selection["page_number"] - 1]]
        selection_x, selection_y, selection_width, selection_height = _selection_bounds(
            pages_to_process[0].rect.width, pages_to_process[0].rect.height, selection
        )
        output_prs = _build_output_presentation(
            pages_to_process[0].rect.width, pages_to_process[0].rect.height
        )
    else:
        pages_to_process = list(doc)
        first_page_rect = doc[0].rect
        selection_x, selection_y, selection_width, selection_height = 0, 0, first_page_rect.width, first_page_rect.height
        output_prs = _build_output_presentation(first_page_rect.width, first_page_rect.height)
    reader = load_ocr()
    report = {
        "pages_processed": len(pages_to_process),
        "pdf_text_blocks": 0,
        "pdf_image_blocks": 0,
        "pdf_pages_ocr": 0,
        "pdf_ocr_regions": 0,
        "pdf_ocr_groups": 0,
        "unsupported": 0,
    }

    for page in pages_to_process:
        new_slide = output_prs.slides.add_slide(output_prs.slide_layouts[6])
        page_selection_rect = _rect_from_xywh(selection_x, selection_y, selection_width, selection_height)
        selection_enabled = selection and selection.get("enabled")

        if selection_enabled:
            _add_picture_from_blob(
                new_slide,
                _render_pdf_page_image(page),
                Pt(0),
                Pt(0),
                Pt(page.rect.width),
                Pt(page.rect.height),
            )

            page_image = _render_pdf_page_for_ocr(page)
            scale_x = page_image.shape[1] / max(page.rect.width, 1)
            scale_y = page_image.shape[0] / max(page.rect.height, 1)
            crop_rect = (
                selection_x * scale_x,
                selection_y * scale_y,
                (selection_x + selection_width) * scale_x,
                (selection_y + selection_height) * scale_y,
            )
            selected_image = _crop_array_to_rect(page_image, crop_rect)
            if selected_image is None:
                report["unsupported"] += 1
                continue
            _overlay_editable_region(
                new_slide,
                selected_image,
                _shape_bounds_object(selection_x, selection_y, selection_width, selection_height),
                report,
                "pdf_ocr",
            )
            report["pdf_pages_ocr"] += 1
            continue

        page_dict = page.get_text("dict")
        text_blocks = _pdf_text_blocks(page_dict)
        image_blocks = _pdf_image_blocks(page_dict)
        page_text_count = 0

        for block in text_blocks:
            if _add_pdf_text_block(new_slide, block, selection_x, selection_y, page_selection_rect):
                page_text_count += 1

        for block in image_blocks:
            if _add_pdf_image_block(new_slide, block, selection_x, selection_y, page_selection_rect):
                report["pdf_image_blocks"] += 1
            else:
                report["unsupported"] += 1

        if page_text_count == 0:
            report["pdf_pages_ocr"] += 1
            ocr_regions, ocr_groups = _ocr_pdf_page(
                page,
                new_slide,
                reader,
                page_selection_rect if selection and selection.get("enabled") else None,
                selection_x,
                selection_y,
            )
            report["pdf_ocr_regions"] += ocr_regions
            report["pdf_ocr_groups"] += ocr_groups
        else:
            report["pdf_text_blocks"] += page_text_count

    return _save_presentation(output_prs), report


def process_image_advanced(input_file, selection=None):
    image_bytes = input_file.read()
    image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    image_np = np.array(image)
    if selection and selection.get("enabled"):
        selection_x, selection_y, selection_width, selection_height = _selection_bounds(
            image.width, image.height, selection
        )
        cropped_np = _crop_array_to_rect(
            image_np,
            (selection_x, selection_y, selection_x + selection_width, selection_y + selection_height),
        )
        if cropped_np is None:
            raise ValueError("The selected image region is empty. Please choose a larger area.")
        slide_width_pts, slide_height_pts = _fit_image_to_slide(image.width, image.height)
    else:
        slide_width_pts, slide_height_pts = _fit_image_to_slide(image.width, image.height)
    output_prs = _build_output_presentation(slide_width_pts, slide_height_pts)
    report = {
        "images_processed": 1,
        "image_ocr_regions": 0,
        "image_text_groups": 0,
    }
    if selection and selection.get("enabled"):
        new_slide = output_prs.slides.add_slide(output_prs.slide_layouts[6])
        full_image_bytes = io.BytesIO()
        image.save(full_image_bytes, format="PNG")
        full_image_bytes.seek(0)
        _add_picture_from_blob(
            new_slide,
            full_image_bytes.getvalue(),
            Pt(0),
            Pt(0),
            output_prs.slide_width,
            output_prs.slide_height,
        )
        left_pt = (selection_x / max(image.width, 1)) * output_prs.slide_width.pt
        top_pt = (selection_y / max(image.height, 1)) * output_prs.slide_height.pt
        width_pt = (selection_width / max(image.width, 1)) * output_prs.slide_width.pt
        height_pt = (selection_height / max(image.height, 1)) * output_prs.slide_height.pt
        _overlay_editable_region(
            new_slide,
            cropped_np,
            _shape_bounds_object(left_pt, top_pt, width_pt, height_pt),
            report,
            "image",
        )
    else:
        _process_flat_image(image_np, output_prs, report, "image")
    return _save_presentation(output_prs), report


def process_uploaded_file(uploaded_file, selection=None):
    extension = Path(uploaded_file.name).suffix.lower()

    if extension == ".pptx":
        result, report = process_pptx_advanced(uploaded_file, selection)
        return result, report, "pptx"

    if extension == ".pdf":
        result, report = process_pdf_advanced(uploaded_file, selection)
        return result, report, "pdf"

    if extension in {".png", ".jpg", ".jpeg"}:
        result, report = process_image_advanced(uploaded_file, selection)
        return result, report, "image"

    raise ValueError(f"Unsupported file type: {extension}")


def _file_summary(uploaded_file):
    extension = Path(uploaded_file.name).suffix.lower().lstrip(".").upper()
    size_mb = uploaded_file.size / (1024 * 1024)
    return extension, f"{size_mb:.2f} MB"


def _uploaded_file_page_count(uploaded_file):
    extension = Path(uploaded_file.name).suffix.lower()
    data = uploaded_file.getvalue()

    if extension == ".pptx":
        return len(Presentation(io.BytesIO(data)).slides)

    if extension == ".pdf" and fitz is not None:
        return fitz.open(stream=data, filetype="pdf").page_count

    return 1


def _selection_preset_values(preset):
    presets = {
        "Full canvas": (0, 0, 100, 100),
        "Top half": (0, 0, 100, 50),
        "Bottom half": (0, 50, 100, 50),
        "Left half": (0, 0, 50, 100),
        "Right half": (50, 0, 50, 100),
        "Center": (20, 20, 60, 60),
    }
    return presets.get(preset)


def _preview_image_for_selection(uploaded_file, selection):
    extension = Path(uploaded_file.name).suffix.lower()
    data = uploaded_file.getvalue()

    if extension in {".png", ".jpg", ".jpeg"}:
        return Image.open(io.BytesIO(data)).convert("RGB")

    if extension == ".pdf" and fitz is not None:
        doc = fitz.open(stream=data, filetype="pdf")
        page = doc[(selection.get("page_number", 1) or 1) - 1]
        pix = page.get_pixmap(matrix=fitz.Matrix(0.4, 0.4), alpha=False)
        return Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")

    if extension == ".pptx":
        prs = Presentation(io.BytesIO(data))
        slide = prs.slides[(selection.get("page_number", 1) or 1) - 1]
        picture_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if len(picture_shapes) == 1:
            return Image.open(io.BytesIO(picture_shapes[0].image.blob)).convert("RGB")

    return None


def _selection_preview_with_overlay(preview_image, selection):
    preview = preview_image.copy()
    draw = ImageDraw.Draw(preview)
    x, y, width, height = _selection_bounds(preview.width, preview.height, selection)
    rect = [x, y, x + width, y + height]
    draw.rectangle(rect, outline=(220, 38, 38), width=max(2, preview.width // 200))
    return preview


def _selection_controls(uploaded_file):
    selection = {"enabled": False}
    with st.expander("Selective Extraction", expanded=False):
        selection["enabled"] = st.checkbox(
            "Extract only a selected section",
            help="Limit extraction to a single region and rebuild it in place on the original canvas.",
        )

        if not selection["enabled"]:
            st.caption("Leave this off to extract the full upload.")
            return selection

        extension = Path(uploaded_file.name).suffix.lower()
        page_count = _uploaded_file_page_count(uploaded_file)
        if extension in {".pptx", ".pdf"}:
            label = "Slide" if extension == ".pptx" else "Page"
            selection["page_number"] = st.number_input(
                f"{label} number",
                min_value=1,
                max_value=page_count,
                value=1,
                step=1,
            )
        else:
            selection["page_number"] = 1

        preset = st.selectbox(
            "Selection preset",
            ["Custom", "Full canvas", "Top half", "Bottom half", "Left half", "Right half", "Center"],
        )
        preset_values = _selection_preset_values(preset)

        if preset_values is None:
            selection["x_pct"] = st.slider("Left (%)", min_value=0, max_value=95, value=0, step=1)
            selection["y_pct"] = st.slider("Top (%)", min_value=0, max_value=95, value=0, step=1)
            selection["width_pct"] = st.slider(
                "Width (%)",
                min_value=1,
                max_value=max(1, 100 - selection["x_pct"]),
                value=max(1, 100 - selection["x_pct"]),
                step=1,
            )
            selection["height_pct"] = st.slider(
                "Height (%)",
                min_value=1,
                max_value=max(1, 100 - selection["y_pct"]),
                value=max(1, 100 - selection["y_pct"]),
                step=1,
            )
        else:
            (
                selection["x_pct"],
                selection["y_pct"],
                selection["width_pct"],
                selection["height_pct"],
            ) = preset_values
            st.caption(
                f"Preset region: left {selection['x_pct']}%, top {selection['y_pct']}%, "
                f"width {selection['width_pct']}%, height {selection['height_pct']}%."
            )

        preview_image = _preview_image_for_selection(uploaded_file, selection)
        if preview_image is not None:
            st.image(
                _selection_preview_with_overlay(preview_image, selection),
                caption="Selection preview",
                use_container_width=True,
            )
        else:
            st.caption(
                "Preview is currently available for images, PDFs, and PowerPoint slides that contain a single picture."
            )

        st.caption(
            "The selected region is defined as a percentage of the current slide, page, or image and is rebuilt "
            "back into the original layout at the same position."
        )
    return selection


def _render_summary_metrics(report, mode):
    if mode == "pptx":
        col1, col2, col3, col4 = st.columns(4)
        col1.metric("Slides", report["slides_processed"])
        col2.metric("OCR Regions", report["pictures_ocr_regions"])
        col3.metric("Grouped Text Boxes", report["pictures_text_groups"])
        col4.metric("Native Shapes", report["text_shapes"] + report["tables"] + report["charts"])
        return

    if mode == "image":
        col1, col2, col3 = st.columns(3)
        col1.metric("Images", report["images_processed"])
        col2.metric("OCR Regions", report["image_ocr_regions"])
        col3.metric("Grouped Text Boxes", report["image_text_groups"])
        return

    col1, col2, col3, col4 = st.columns(4)
    col1.metric("Pages", report["pages_processed"])
    col2.metric("Native Text Blocks", report["pdf_text_blocks"])
    col3.metric("Image Blocks", report["pdf_image_blocks"])
    col4.metric("OCR Pages", report["pdf_pages_ocr"])


def _render_report_highlights(report, mode):
    if mode == "pptx":
        if report["pictures_ocr_regions"] > 0:
            st.info(
                f"Detected {report['pictures_ocr_regions']} OCR text region(s) inside flattened slide images "
                f"and merged them into {report['pictures_text_groups']} larger editable text box(es)."
            )
        if report["text_shapes"] + report["tables"] + report["charts"] > 0:
            st.success(
                f"Preserved {report['text_shapes']} text/shape elements, "
                f"{report['tables']} tables, and {report['charts']} charts as native PowerPoint objects."
            )
        return

    if mode == "image":
        st.success(
            f"Grouped {report['image_ocr_regions']} OCR region(s) into {report['image_text_groups']} editable "
            "text box(es) on a single-slide PowerPoint output."
        )
        return

    if report["pdf_pages_ocr"] > 0:
        st.warning(
            f"{report['pdf_pages_ocr']} PDF page(s) had no native text layer and were processed with OCR fallback."
        )
        if report["pdf_ocr_regions"] > 0:
            st.info(
                f"OCR found {report['pdf_ocr_regions']} text region(s) and grouped them into "
                f"{report['pdf_ocr_groups']} larger text box(es)."
            )
    else:
        st.success("All PDF pages were reconstructed from native PDF text and image blocks.")

    if report["pdf_text_blocks"] > 0:
        st.info(f"Recovered {report['pdf_text_blocks']} native text block(s) as editable text boxes.")


st.set_page_config(page_title="Deck And PDF Reconstructor", layout="wide")

with st.sidebar:
    st.subheader("What This Does")
    st.write(
        "Rebuilds uploaded decks and PDFs into editable PowerPoint output. Native elements stay editable "
        "when possible, and OCR is used only when content is flattened."
    )
    st.subheader("Supported Inputs")
    st.write("`PPTX`: native shapes, tables, charts, and screenshot-based OCR recovery")
    st.write("`PDF`: native text blocks, embedded images, and OCR fallback for scanned pages")
    st.write("`PNG/JPG`: single-image OCR reconstruction into an editable one-slide PPTX")
    st.subheader("Current Limits")
    st.write("Complex vector graphics, SmartArt, and arbitrary PDF drawing paths are still best-effort.")
    st.write("Selective extraction currently works best on image-based content and screenshot-style slides.")

st.title("Deck And PDF Reconstructor")
st.caption(
    "Upload a PowerPoint deck, PDF, or image to reconstruct editable text, preserve native elements, "
    "and separate flattened content into reusable pieces."
)

uploaded_file = st.file_uploader(
    "Upload a `.pptx`, `.pdf`, `.png`, `.jpg`, or `.jpeg` file",
    type=["pptx", "pdf", "png", "jpg", "jpeg"],
)

if uploaded_file:
    extension, size_label = _file_summary(uploaded_file)
    selection = _selection_controls(uploaded_file)
    preview_col1, preview_col2, preview_col3 = st.columns([1, 1, 2])
    preview_col1.metric("Format", extension)
    preview_col2.metric("Size", size_label)
    preview_col3.info(
        "Processing mode: native reconstruction first, with OCR regions merged into larger grouped text boxes."
    )

    if st.button("Extract Everything", type="primary", use_container_width=True):
        with st.spinner("Reconstructing editable content..."):
            try:
                uploaded_file.seek(0)
                result, report, mode = process_uploaded_file(uploaded_file, selection)
            except Exception as exc:
                st.exception(exc)
            else:
                output_name = f"{Path(uploaded_file.name).stem}_editable.pptx"
                st.success("Reconstruction complete.")
                _render_summary_metrics(report, mode)
                _render_report_highlights(report, mode)

                details_col1, details_col2 = st.columns([3, 2])
                with details_col1:
                    st.subheader("Extraction Report")
                    st.json(report)
                with details_col2:
                    st.subheader("Output")
                    st.write("Your editable reconstruction is ready as a PowerPoint file.")
                    st.download_button(
                        "Download reconstructed PPTX",
                        result,
                        output_name,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                    )
else:
    st.info(
        "Drop in a PowerPoint deck, PDF, or image to see the extraction summary and download an editable PPTX."
    )
